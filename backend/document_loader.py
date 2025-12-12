"""
Document loader for Notion Markdown and PowerPoint files.
Handles text extraction from various document formats.
"""

import os
from pathlib import Path
from typing import List, Dict, Optional
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation


class DocumentLoader:
    """Loads and extracts text from Notion Markdown and PowerPoint files."""
    
    def __init__(self):
        self.supported_extensions = {'.md', '.markdown', '.pptx'}
    
    def load_document(self, file_path: str) -> Dict[str, any]:
        """
        Load a document and extract its content.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary with 'content', 'metadata', and 'type' keys
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = path.suffix.lower()
        
        if extension == '.md' or extension == '.markdown':
            return self._load_markdown(file_path)
        elif extension == '.pptx':
            return self._load_powerpoint(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")
    
    def _load_markdown(self, file_path: str) -> Dict[str, any]:
        """Load and parse Notion Markdown file."""
        path = Path(file_path)
        
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Parse markdown to extract text
        html = markdown.markdown(content, extensions=['tables', 'fenced_code'])
        soup = BeautifulSoup(html, 'html.parser')
        
        # Extract text content
        text_content = soup.get_text(separator='\n', strip=True)
        
        # Extract title (first H1 or filename)
        title = path.stem
        h1_tags = soup.find_all('h1')
        if h1_tags:
            title = h1_tags[0].get_text(strip=True)
        
        # Extract metadata
        metadata = {
            'title': title,
            'filename': path.name,
            'file_path': str(path),
            'file_type': 'markdown',
            'word_count': len(text_content.split())
        }
        
        return {
            'content': text_content,
            'raw_content': content,  # Keep raw markdown for better chunking
            'metadata': metadata,
            'type': 'markdown'
        }
    
    def _load_powerpoint(self, file_path: str) -> Dict[str, any]:
        """Load and parse PowerPoint file."""
        path = Path(file_path)
        prs = Presentation(file_path)
        
        slides_content = []
        slide_metadata = []
        
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # Extract notes if available
            notes_text = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
            
            # Combine slide text and notes
            full_slide_text = "\n".join(slide_text)
            if notes_text:
                full_slide_text += f"\n\n[Notes: {notes_text}]"
            
            if full_slide_text.strip():
                slides_content.append(full_slide_text)
                slide_metadata.append({
                    'slide_number': idx,
                    'title': slide_text[0] if slide_text else f"Slide {idx}",
                    'has_notes': bool(notes_text)
                })
        
        # Combine all slides
        full_content = "\n\n---\n\n".join([
            f"Slide {i+1}:\n{content}" 
            for i, content in enumerate(slides_content)
        ])
        
        metadata = {
            'title': path.stem,
            'filename': path.name,
            'file_path': str(path),
            'file_type': 'powerpoint',
            'total_slides': len(slides_content),
            'slide_metadata': slide_metadata,
            'word_count': len(full_content.split())
        }
        
        return {
            'content': full_content,
            'raw_content': full_content,
            'metadata': metadata,
            'type': 'powerpoint',
            'slides': slides_content  # Keep individual slides for chunking
        }
    
    def load_directory(self, directory_path: str) -> List[Dict[str, any]]:
        """
        Load all supported documents from a directory.
        
        Args:
            directory_path: Path to directory containing documents
            
        Returns:
            List of document dictionaries
        """
        path = Path(directory_path)
        if not path.is_dir():
            raise ValueError(f"Not a directory: {directory_path}")
        
        documents = []
        for file_path in path.iterdir():
            if file_path.suffix.lower() in self.supported_extensions:
                try:
                    doc = self.load_document(str(file_path))
                    documents.append(doc)
                except Exception as e:
                    print(f"Error loading {file_path}: {e}")
                    continue
        
        return documents

